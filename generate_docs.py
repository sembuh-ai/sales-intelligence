#!/usr/bin/env python3
"""
Sales Intelligence — Document Generator
Fetches deals from Monday.com, generates Quotation/Pricing/Proposal docs,
uploads to Google Drive, creates Gmail draft, notifies Slack.
"""

import copy
import json
import os
import re
import sys
import datetime
import requests

from dotenv import load_dotenv

# Load all .env files (monday/ has Slack tokens, gmcp/ has Google config)
_base = os.path.dirname(os.path.abspath(__file__))
load_dotenv(os.path.join(_base, ".env"))
load_dotenv(os.path.join(_base, "monday", ".env"), override=False)
load_dotenv(os.path.join(_base, "gmcp", ".env"), override=False)

# ── Config ────────────────────────────────────────────────────
MONDAY_API_KEY = os.getenv("MONDAY_API_KEY", "")
MONDAY_WORKSPACE_ID = os.getenv("MONDAY_WORKSPACE_ID", "")
SLACK_BOT_TOKEN = os.getenv("SLACK_BOT_TOKEN", "")
SLACK_CHANNEL_ID = os.getenv("SLACK_CHANNEL_ID", "C0AU9FA76EB")
GOOGLE_FOLDER_ID = os.getenv("FOLDER_ID", "")

MONDAY_API_URL = "https://api.monday.com/v2"
SLACK_API_URL = "https://slack.com/api"

TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "PTP Hackathon - Brief")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "output")

TODAY = datetime.date.today()
TODAY_STR = TODAY.strftime("%d %B %Y")
TODAY_ISO = TODAY.isoformat()

# Pricing assumptions (from template)
MARGIN_ASSUMPTIONS = {
    "Claim Workflow": {"unit_price": 1.0, "margin": 0.60, "type": "Recurring", "uom": "Per Claim"},
    "Fraud Detection": {"unit_price": 2.0, "margin": 0.65, "type": "Recurring", "uom": "Per Claim"},
    "Implementation": {"unit_price": 50000, "margin": 0.50, "type": "One Time Fee", "uom": "Per Implementation"},
}

# Stage weights for revenue forecast
STAGE_WEIGHTS = {
    "Open": 0.10,
    "First Meeting Done": 0.20,
    "Piloting": 0.30,
    "Solutioning": 0.50,
    "Waiting Confirmation": 0.70,
    "[Won] Waiting Signature": 0.90,
    "Implementation Done": 1.00,
    "Lost": 0.00,
}

# Product mapping based on deal name keywords
def detect_products(deal_name):
    """Detect Sembuh AI products from deal name."""
    name_lower = deal_name.lower()
    products = []
    if "full suite" in name_lower:
        products = ["Claim Workflow", "Fraud Detection", "Implementation"]
    elif "fwa" in name_lower or "fraud" in name_lower:
        products = ["Fraud Detection", "Implementation"]
    elif "ocr" in name_lower or "stp" in name_lower:
        products = ["Claim Workflow", "Implementation"]
    else:
        products = ["Claim Workflow", "Fraud Detection", "Implementation"]
    return products


# ── Monday.com API ────────────────────────────────────────────

def monday_query(query, variables=None):
    """Execute Monday.com GraphQL query."""
    headers = {
        "Authorization": MONDAY_API_KEY,
        "Content-Type": "application/json",
        "API-Version": "2024-10",
    }
    payload = {"query": query}
    if variables:
        payload["variables"] = variables
    resp = requests.post(MONDAY_API_URL, json=payload, headers=headers)
    resp.raise_for_status()
    data = resp.json()
    if "errors" in data:
        raise RuntimeError(f"Monday API error: {data['errors']}")
    return data["data"]


def fetch_boards():
    """Fetch boards in workspace."""
    query = """
    query ($wsIds: [ID!]) {
        boards(workspace_ids: $wsIds, limit: 50) {
            id
            name
        }
    }
    """
    data = monday_query(query, {"wsIds": [int(MONDAY_WORKSPACE_ID)]})
    return data["boards"]


def fetch_deals(board_id):
    """Fetch all items (deals) from a board with column values."""
    query = """
    query ($boardId: [ID!]!) {
        boards(ids: $boardId) {
            items_page(limit: 100) {
                items {
                    id
                    name
                    column_values {
                        id
                        type
                        text
                        value
                    }
                }
            }
        }
    }
    """
    data = monday_query(query, {"boardId": [int(board_id)]})
    items = data["boards"][0]["items_page"]["items"]
    deals = []
    for item in items:
        deal = {"id": item["id"], "name": item["name"]}
        for col in item["column_values"]:
            col_id = col["id"]
            text = col.get("text", "") or ""
            raw = col.get("value", "")
            deal[col_id] = text
            deal[f"{col_id}_raw"] = raw
        deals.append(deal)
    return deals


def parse_deal(deal):
    """Parse raw deal dict into structured data."""
    name = deal.get("name", "Unknown")
    stage = deal.get("deal_stage", "Open")
    deal_value = _parse_num(deal.get("deal_value", "0"))
    close_date = deal.get("deal_expected_close_date", "")
    monthly_claim_vol = _parse_num(deal.get("numeric_mm1bmx9t", "0"))
    members_covered = _parse_num(deal.get("numeric_mm1bx91m", "0"))
    pricing_model = deal.get("dropdown_mm1b79r5", "Per Claim")
    incurred = _parse_num(deal.get("numeric_mm1bdpzy", "0"))
    excess = _parse_num(deal.get("numeric_mm1bkxy8", "0"))
    approved = _parse_num(deal.get("numeric_mm1b64b7", "0"))
    op_pct = _parse_num(deal.get("numeric_mm1bq8yc", "0"))
    ip_pct = _parse_num(deal.get("numeric_mm1bek79", "0"))
    owner = deal.get("deal_owner", "")
    proposal_date = deal.get("date_mm1bpvvx", "")

    products = detect_products(name)

    # Estimate volume from deal_value if not set
    if monthly_claim_vol == 0 and deal_value > 0:
        # Reverse-engineer from deal_value: assume ~$3/claim avg (CW $1 + FWA $2) + impl fee
        impl_fee = 50000 if "Implementation" in products else 0
        recurring_rate = sum(
            MARGIN_ASSUMPTIONS[p]["unit_price"] for p in products if p != "Implementation"
        ) or 1
        annual_claim_vol = max(int((deal_value - impl_fee) / recurring_rate), 1200) if deal_value > impl_fee else 12000
        monthly_claim_vol = annual_claim_vol / 12
    else:
        annual_claim_vol = monthly_claim_vol * 12

    if members_covered == 0 and monthly_claim_vol > 0:
        members_covered = monthly_claim_vol * 10  # rough estimate

    return {
        "id": deal.get("id"),
        "name": name,
        "stage": stage,
        "deal_value": deal_value,
        "close_date": close_date,
        "monthly_claim_vol": monthly_claim_vol,
        "annual_claim_vol": annual_claim_vol,
        "members_covered": members_covered,
        "pricing_model": pricing_model or "Per Claim",
        "incurred": incurred,
        "excess": excess,
        "approved": approved,
        "op_pct": op_pct,
        "ip_pct": ip_pct,
        "owner": owner,
        "proposal_date": proposal_date,
        "products": products,
    }


def _parse_num(val):
    if not val:
        return 0
    val = str(val).replace(",", "").replace("$", "").strip()
    try:
        return float(val)
    except ValueError:
        return 0


# ── Quotation Generator ──────────────────────────────────────

def generate_quotation(deal, seq_num=1):
    """Generate quotation xlsx from template, filled with deal data."""
    import openpyxl

    template_path = os.path.join(TEMPLATE_DIR, "Sembuh AI_Quotation Template.xlsx")
    wb = openpyxl.load_workbook(template_path)
    ws = wb.active

    client_name = deal["name"]
    q_number = f"SBH-Q-{TODAY.year}-{seq_num:03d}"
    validity_date = (TODAY + datetime.timedelta(days=30)).strftime("%d %B %Y")

    products = deal["products"]
    products_str = ", ".join(products)
    monthly_vol = int(deal["monthly_claim_vol"] or 0)
    annual_vol = int(deal["annual_claim_vol"] or monthly_vol * 12)
    members = int(deal["members_covered"] or monthly_vol * 10)

    # Compute PMPM cost
    recurring_rate = sum(MARGIN_ASSUMPTIONS[p]["unit_price"] for p in products if p != "Implementation")
    impl_fee = MARGIN_ASSUMPTIONS["Implementation"]["unit_price"] if "Implementation" in products else 0
    total_annual = recurring_rate * annual_vol + impl_fee
    pmpm = total_annual / members / 12 if members > 0 else 0

    # Replace placeholders in-place (preserve formatting, merges, layout)
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
        for cell in row:
            if cell.value is None:
                continue
            val = str(cell.value)
            if "[Placeholder" not in val and "Placeholder" not in val and "~900,000" not in val and "~72,000" not in val:
                continue
            val = val.replace("[Placeholder Name]", deal.get("owner", "Account Manager"))
            val = val.replace("[Placeholder Client]", client_name)
            val = val.replace("[Placeholder Full Address]", "")
            val = val.replace("Placeholder DD/Month/YYYY", TODAY_STR)
            val = val.replace("[Placeholder DD/Month/YYYY]", TODAY_STR)
            if val.strip() in ("[Placeholder]", "[Placeholder] "):
                if cell.row == 10:
                    val = f"{products_str} for {client_name}"
                elif cell.column == 9:  # Cost column (I)
                    val = f"IDR {pmpm * 16000:,.0f}" if pmpm > 0 else "Per Claim"
                else:
                    val = client_name
            val = val.replace("~900,000 members", f"~{members:,} members")
            val = val.replace("~72,000 claims", f"~{annual_vol:,} claims")
            cell.value = val

    # Fill header fields
    ws["H4"] = f"Quotation No: {q_number}"
    ws["H5"] = f"Date: {TODAY_STR}"
    ws["H6"] = f"Valid Until: {validity_date}"

    safe_name = re.sub(r'[^\w\s-]', '', client_name).strip().replace(' ', '_')
    out_path = os.path.join(OUTPUT_DIR, safe_name, f"Sembuh AI - {client_name} - Quotation v1.xlsx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    wb.save(out_path)
    print(f"  Quotation saved: {out_path}")
    return out_path


# ── Pricing Internal Generator ────────────────────────────────

def generate_pricing_internal(deal, all_deals):
    """Generate pricing internal xlsx from template, filled with deal data."""
    import openpyxl

    template_path = os.path.join(TEMPLATE_DIR, "Hackathon_Pricing Internal.xlsx")
    wb = openpyxl.load_workbook(template_path)

    # --- Sheet 1: Per-Client (named after client) ---
    client_name = deal["name"]
    products = deal["products"]

    # Use first sheet as per-client sheet
    ws1 = wb.worksheets[0]
    ws1.title = client_name[:31]  # Excel sheet name max 31 chars

    # Unmerge and clear
    for merge in list(ws1.merged_cells.ranges):
        ws1.unmerge_cells(str(merge))
    for row in ws1.iter_rows(min_row=1, max_row=ws1.max_row, max_col=ws1.max_column):
        for cell in row:
            cell.value = None

    ws1["A1"] = "Per Client Gross Profit"
    ws1["H1"] = f"Date {TODAY_STR}"
    ws1["A2"] = "Client"
    ws1["B2"] = "Status"
    ws1["C2"] = "Total Revenue ($)"
    ws1["D2"] = "Blended Margin %"
    ws1["E2"] = "Cost of Revenue ($)"
    ws1["F2"] = "Gross Profit ($)"

    # Calculate total revenue for this client
    total_rev = 0
    for prod_name in products:
        info = MARGIN_ASSUMPTIONS[prod_name]
        if prod_name == "Implementation":
            total_rev += info["unit_price"]
        else:
            total_rev += info["unit_price"] * deal["annual_claim_vol"]

    blended_margin = 0.60
    if products:
        margins = [MARGIN_ASSUMPTIONS[p]["margin"] for p in products]
        blended_margin = sum(margins) / len(margins)

    ws1["A3"] = client_name
    ws1["B3"] = deal["stage"]
    ws1["C3"] = round(total_rev, 2)
    ws1["D3"] = round(blended_margin, 2)
    ws1["E3"] = f"=C3*(1-D3)"
    ws1["F3"] = f"=C3*D3"

    # --- Sheet 2: Margin Summary (update with all deals) ---
    ws2 = wb.worksheets[1] if len(wb.worksheets) > 1 else wb.create_sheet("Margin Summary")
    for merge in list(ws2.merged_cells.ranges):
        ws2.unmerge_cells(str(merge))
    for row in ws2.iter_rows(min_row=1, max_row=ws2.max_row, max_col=ws2.max_column):
        for cell in row:
            cell.value = None

    ws2["A1"] = "Sembuh AI — Gross Profit & Margin Summary"
    ws2["A2"] = "Margin assumptions: Claim Workflow 60% | Fraud Detection 65% | Implementation 50%"

    ws2["A4"] = "Product / Service"
    ws2["B4"] = "Type"
    ws2["C4"] = "Unit Price ($)"
    ws2["D4"] = "Margin %"
    ws2["E4"] = "Cost per Unit ($)"
    ws2["F4"] = "Gross Profit per Unit ($)"
    ws2["G4"] = "Notes"

    margin_data = [
        ("Claim Workflow", "Recurring", 1.0, 0.60, "AI-enabled claim automation — low infra cost per claim"),
        ("Fraud Detection", "Recurring", 2.0, 0.65, "ML inference cost is low relative to price point"),
        ("Implementation", "One Time Fee", 50000, 0.50, "Engineering + integration effort — 50% margin floor"),
    ]
    for i, (name, typ, price, margin, notes) in enumerate(margin_data):
        r = 5 + i
        ws2[f"A{r}"] = name
        ws2[f"B{r}"] = typ
        ws2[f"C{r}"] = price
        ws2[f"D{r}"] = margin
        ws2[f"E{r}"] = f"=C{r}*(1-D{r})"
        ws2[f"F{r}"] = f"=C{r}*D{r}"
        ws2[f"G{r}"] = notes

    # Per-client summary
    ws2["A10"] = "Per Client Gross Profit"
    ws2["A11"] = "Client"
    ws2["B11"] = "Status"
    ws2["C11"] = "Total Revenue ($)"
    ws2["D11"] = "Blended Margin %"
    ws2["E11"] = "Cost of Revenue ($)"
    ws2["F11"] = "Gross Profit ($)"
    ws2["G11"] = "GP as % of Pipeline"

    for i, d in enumerate(all_deals):
        r = 12 + i
        prods = d["products"]
        rev = 0
        for p in prods:
            info = MARGIN_ASSUMPTIONS[p]
            if p == "Implementation":
                rev += info["unit_price"]
            else:
                rev += info["unit_price"] * d["annual_claim_vol"]
        margins = [MARGIN_ASSUMPTIONS[p]["margin"] for p in prods] if prods else [0.6]
        bm = sum(margins) / len(margins)

        ws2[f"A{r}"] = d["name"]
        ws2[f"B{r}"] = d["stage"]
        ws2[f"C{r}"] = round(rev, 2)
        ws2[f"D{r}"] = round(bm, 2)
        ws2[f"E{r}"] = f"=C{r}*(1-D{r})"
        ws2[f"F{r}"] = f"=C{r}*D{r}"

    total_r = 12 + len(all_deals)
    ws2[f"A{total_r}"] = "TOTAL"
    ws2[f"C{total_r}"] = f"=SUM(C12:C{total_r - 1})"
    ws2[f"D{total_r}"] = f"=F{total_r}/C{total_r}"
    ws2[f"E{total_r}"] = f"=SUM(E12:E{total_r - 1})"
    ws2[f"F{total_r}"] = f"=SUM(F12:F{total_r - 1})"

    # --- Sheet 3: Deal Summary ---
    ws3 = wb.worksheets[2] if len(wb.worksheets) > 2 else wb.create_sheet("Deal Summary")
    for merge in list(ws3.merged_cells.ranges):
        ws3.unmerge_cells(str(merge))
    for row in ws3.iter_rows(min_row=1, max_row=ws3.max_row, max_col=ws3.max_column):
        for cell in row:
            cell.value = None

    ws3["A1"] = "Sembuh AI — Potential Client Pipeline"
    ws3["A2"] = "Products & Services Proposal Summary  |  Annual Basis  |  IDR pricing converted to USD at 16,000"

    headers = ["Client", "Status", "Products Purchased", "Monthly Claim Vol", "Annual Claim Vol",
               "Claim Workflow\n($/claim)", "Fraud Detection\n($/claim)", "Claim Workflow\nRevenue ($)",
               "Fraud Detection\nRevenue ($)", "Implementation\nFee ($)", "Total Annual\nRevenue ($)",
               "Notes", "Blended\nMargin %", "Cost of\nRevenue ($)", "Gross\nProfit ($)"]
    for i, h in enumerate(headers):
        ws3.cell(row=4, column=i + 1, value=h)

    for i, d in enumerate(all_deals):
        r = 5 + i
        prods = d["products"]
        has_cw = "Claim Workflow" in prods
        has_fd = "Fraud Detection" in prods
        has_impl = "Implementation" in prods
        cw_price = 1.0 if has_cw else 0
        fd_price = 2.0 if has_fd else 0
        impl_fee = 50000 if has_impl else 0

        margins = [MARGIN_ASSUMPTIONS[p]["margin"] for p in prods] if prods else [0.6]
        bm = sum(margins) / len(margins)

        ws3[f"A{r}"] = d["name"]
        ws3[f"B{r}"] = d["stage"]
        ws3[f"C{r}"] = " + ".join(prods)
        ws3[f"D{r}"] = d["monthly_claim_vol"]
        ws3[f"E{r}"] = f"=D{r}*12"
        ws3[f"F{r}"] = cw_price if has_cw else "-"
        ws3[f"G{r}"] = fd_price if has_fd else "-"
        ws3[f"H{r}"] = f"=E{r}*F{r}" if has_cw else ""
        ws3[f"I{r}"] = f"=E{r}*G{r}" if has_fd else ""
        ws3[f"J{r}"] = impl_fee if has_impl else ""
        ws3[f"K{r}"] = f"=H{r}+I{r}+J{r}"
        ws3[f"L{r}"] = ""
        ws3[f"M{r}"] = round(bm, 2)
        ws3[f"N{r}"] = f"=K{r}*(1-M{r})"
        ws3[f"O{r}"] = f"=K{r}*M{r}"

    # --- Sheet 4: Client Detail ---
    ws4 = wb.worksheets[3] if len(wb.worksheets) > 3 else wb.create_sheet("Client Detail")
    for merge in list(ws4.merged_cells.ranges):
        ws4.unmerge_cells(str(merge))
    for row in ws4.iter_rows(min_row=1, max_row=ws4.max_row, max_col=ws4.max_column):
        for cell in row:
            cell.value = None

    ws4["A1"] = "Sembuh AI — Per Client Product & Pricing Detail"
    ws4["A2"] = "Annual revenue breakdown per client  |  Prices in USD  |  Volume = annual claims"

    detail_headers = ["Client", "Product / Service", "Type", "Unit Price ($)", "Annual Volume",
                      "UoM", "Annual Revenue ($)", "Notes", "Margin %", "Cost ($)", "Gross Profit ($)"]
    for i, h in enumerate(detail_headers):
        ws4.cell(row=4, column=i + 1, value=h)

    detail_notes = {
        "Claim Workflow": "AI-enabled end-to-end claim workflow",
        "Fraud Detection": "FWA detection — fraud, waste & abuse",
        "Implementation": "Integration into insurance core system",
    }

    row_num = 5
    for d in all_deals:
        for prod_name in d["products"]:
            info = MARGIN_ASSUMPTIONS[prod_name]
            vol = 1 if prod_name == "Implementation" else d["annual_claim_vol"]
            ws4[f"A{row_num}"] = d["name"]
            ws4[f"B{row_num}"] = prod_name
            ws4[f"C{row_num}"] = info["type"]
            ws4[f"D{row_num}"] = info["unit_price"]
            ws4[f"E{row_num}"] = vol
            ws4[f"F{row_num}"] = info["uom"]
            ws4[f"G{row_num}"] = f"=D{row_num}*E{row_num}"
            ws4[f"H{row_num}"] = detail_notes.get(prod_name, "")
            ws4[f"I{row_num}"] = info["margin"]
            ws4[f"J{row_num}"] = f"=G{row_num}*(1-I{row_num})"
            ws4[f"K{row_num}"] = f"=G{row_num}*I{row_num}"
            row_num += 1

    safe_name = re.sub(r'[^\w\s-]', '', client_name).strip().replace(' ', '_')
    out_path = os.path.join(OUTPUT_DIR, safe_name, f"Sembuh AI - {client_name} - Pricing Internal.xlsx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    wb.save(out_path)
    print(f"  Pricing Internal saved: {out_path}")
    return out_path


# ── Proposal Generator ────────────────────────────────────────

def _compute_proposal_metrics(deal, discount_pct=0):
    """Compute ROI and executive summary metrics from deal data."""
    monthly_vol = deal["monthly_claim_vol"] or 4000
    annual_vol = monthly_vol * 12
    members = deal["members_covered"] or monthly_vol * 10
    incurred = deal["incurred"] or 0
    excess = deal["excess"] or 0
    approved = deal["approved"] or (incurred - excess if incurred else 0)

    # Excess rate
    excess_rate = (excess / incurred * 100) if incurred > 0 else 4.0
    sim_excess_rate = min(excess_rate + 3, 15)  # +3% improvement scenario
    sim_excess = incurred * sim_excess_rate / 100 if incurred > 0 else 0
    containment_value = sim_excess - excess if sim_excess > excess else 0

    # Cost model: per claim pricing (with discount)
    discount_mult = 1 - (discount_pct / 100) if discount_pct else 1.0
    products = deal["products"]
    has_cw = "Claim Workflow" in products
    has_fd = "Fraud Detection" in products
    has_impl = "Implementation" in products
    cw_price = 1.0 * discount_mult if has_cw else 0
    fd_price = 2.0 * discount_mult if has_fd else 0
    impl_fee = 50000 * discount_mult if has_impl else 0
    annual_recurring = (cw_price + fd_price) * annual_vol
    annual_investment = annual_recurring + impl_fee

    # PMPM (per member per month)
    pmpm = annual_investment / members / 12 if members > 0 else 0

    # Hours saved (~4min/claim savings)
    hours_saved = annual_vol * 4 / 60
    productivity_value = hours_saved * 50000  # IDR 50K/hour

    # Operational efficiency
    time_before = 10  # min/claim
    time_after = 6    # min/claim
    pct_faster = round((time_before - time_after) / time_before * 100)

    # ROI
    total_value = containment_value + productivity_value / 16000  # Convert IDR to USD
    roi_multiple = total_value / annual_investment if annual_investment > 0 else 0

    # IDR conversions (assume 16,000 IDR/USD)
    fx = 16000
    incurred_idr = incurred * fx if incurred < 1e9 else incurred
    annual_inv_idr = annual_investment * fx

    return {
        "members": int(members),
        "annual_vol": int(annual_vol),
        "monthly_vol": int(monthly_vol),
        "incurred_idr": incurred_idr,
        "excess_rate": round(excess_rate, 1),
        "sim_excess_rate": round(sim_excess_rate, 1),
        "containment_value": containment_value,
        "annual_investment": annual_investment,
        "annual_inv_idr": annual_inv_idr,
        "pmpm_idr": round(pmpm * fx),
        "hours_saved": int(hours_saved),
        "productivity_idr": int(productivity_value),
        "pct_faster": pct_faster,
        "roi_multiple": round(roi_multiple, 1),
        "total_value": total_value,
        "time_before": time_before,
        "time_after": time_after,
        "products": products,
        "discount_pct": discount_pct,
        "cw_price": cw_price,
        "fd_price": fd_price,
        "impl_fee": impl_fee,
    }


def _fmt_idr(val):
    """Format number as IDR billions/millions."""
    if abs(val) >= 1e9:
        return f"IDR {val/1e9:,.1f} B"
    if abs(val) >= 1e6:
        return f"IDR {val/1e6:,.0f} M"
    return f"IDR {val:,.0f}"


def generate_proposal(deal, discount_pct=0, discount_note="", mom_context=""):
    """Generate proposal pptx from template, replacing placeholders with deal data."""
    from pptx import Presentation

    template_path = os.path.join(TEMPLATE_DIR, "Sembuh AI_Proposal.pptx")
    prs = Presentation(template_path)

    client_name = deal["name"]
    products_str = ", ".join(deal["products"])
    m = _compute_proposal_metrics(deal, discount_pct=discount_pct)

    # --- Build replacement map for text across all slides ---
    replacements = {
        # Date placeholders
        "[Placeholder DD/Month/YYYY]": TODAY_STR,
        "DD/Month/YYYY": TODAY_STR,
        "Placeholder DD/Month/YYYY": TODAY_STR,
        # Client name (BNI Life is the template default)
        "BNI Life": client_name,
        "BNI LIfe": client_name,
        "BNI LIFE": client_name.upper(),
    }

    # --- Key Executive Summary placeholders (Slide 3, shape S85) ---
    scope_text = (
        f"Sembuh AI provides {', '.join(deal['products']).lower()} "
        f"supporting both API and web-app workflows."
    )
    if "Claim Workflow" in deal["products"] and "Fraud Detection" in deal["products"]:
        scope_text = (
            "Sembuh AI provides client-hosted claim data extraction and "
            "Fraud, Waste & Abuse detection, supporting both API and web-app workflows."
        )

    context_text = (
        f"{client_name} aims to improve speed, consistency, and auditability of "
        f"claims processing through AI-powered OCR extraction, workflow automation, "
        f"and intelligent cost containment analysis."
    )

    discount_str = ""
    if discount_pct:
        discount_str = f" (includes {discount_pct}% discount as agreed)"

    benchmark_text = (
        f"Coverage: {m['members']:,} members, ~{m['annual_vol']:,} claims/year. "
        f"Commercial model: per-claim pricing{discount_str}. "
        f"Estimated annual investment: {_fmt_idr(m['annual_inv_idr'])}."
    )

    if mom_context:
        approach_text = (
            f"Based on our recent discussions with {client_name}: {mom_context} "
            f"Implementation timeline to be finalized after IT alignment."
        )
    else:
        approach_text = (
            f"Structured PoC using {client_name} claim samples to validate: "
            f"(1) OCR extraction quality, (2) Claim Workflow Automation & FWA analysis alignment, "
            f"and (3) workflow fit. Implementation timeline to be finalized after IT alignment."
        )

    placeholder_replacements = {
        # Shape S85 [Placeholder] markers
        "Context & Objective": "Context & Objective",
        "Scope of Solution ": "Scope of Solution ",
        "Approach & Moving Forward": "Approach & Moving Forward",
    }

    # --- Slide 14: Key Executive Summary ROI numbers ---
    pmpm_label = f"IDR {m['pmpm_idr']:,} per member per month"
    if discount_pct:
        pmpm_label += f" ({discount_pct}% discount applied)"

    roi_replacements = {
        "~ IDR 61.94B+": f"~ {_fmt_idr(m['total_value'] * 16000)}+",
        "~ IDR 6.934 B": f"~ {_fmt_idr(m['annual_inv_idr'])}",
        "~ 9 x ": f"~ {m['roi_multiple']:.0f} x " if m['roi_multiple'] >= 1 else "< 1 x ",
        "900,000 members and ~72,000 claims/year": f"{m['members']:,} members and ~{m['annual_vol']:,} claims/year",
        "IDR 650 per member per month [1]": pmpm_label,
        "CNI-QT-26-03-00070": f"Based on {client_name} deal data",
    }
    replacements.update(roi_replacements)

    # --- Slide 15: ROI excess scenario ---
    excess_replacements = {
        "IDR 2.310 B": _fmt_idr(m["incurred_idr"]) if m["incurred_idr"] > 0 else "IDR 2.310 B",
        "IDR 2.210 B": _fmt_idr(m["incurred_idr"] * (100 - m["excess_rate"]) / 100) if m["incurred_idr"] > 0 else "IDR 2.210 B",
        "IDR 100 B": _fmt_idr(m["incurred_idr"] * m["excess_rate"] / 100) if m["incurred_idr"] > 0 else "IDR 100 B",
        "IDR 2.148 B": _fmt_idr(m["incurred_idr"] * (100 - m["sim_excess_rate"]) / 100) if m["incurred_idr"] > 0 else "IDR 2.148 B",
        "IDR 161.7B": _fmt_idr(m["incurred_idr"] * m["sim_excess_rate"] / 100) if m["incurred_idr"] > 0 else "IDR 161.7B",
    }
    replacements.update(excess_replacements)

    # --- Slide 17: Operational efficiency ---
    ops_replacements = {
        "~4,800": f"~{m['hours_saved']:,}",
        "~72.000/year": f"~{m['annual_vol']:,}/year",
        "~288,000 minutes": f"~{m['annual_vol'] * 4:,} minutes",
        "~4,800 hours": f"~{m['hours_saved']:,} hours",
        "~240M": f"~{m['productivity_idr'] / 1e6:,.0f}M",
        "72,000 claims/year": f"{m['annual_vol']:,} claims/year",
        "4,800 hours saved annually": f"{m['hours_saved']:,} hours saved annually",
        "~IDR 240M/year productivity value": f"~{_fmt_idr(m['productivity_idr'])}/year productivity value",
    }
    replacements.update(ops_replacements)

    # --- Apply all replacements across slides ---
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs_list = list(para.runs)
                    for run_idx, run in enumerate(runs_list):
                        # Replace [Placeholder] markers with dynamic content
                        if run.text.strip() == "[Placeholder]":
                            prev_text = runs_list[run_idx - 1].text if run_idx > 0 else ""
                            if "Context" in prev_text:
                                run.text = context_text
                            elif "Scope" in prev_text:
                                run.text = scope_text
                            elif "Benchmark" in prev_text or "Reference" in prev_text:
                                run.text = benchmark_text
                            elif "Approach" in prev_text or "Moving Forward" in prev_text:
                                run.text = approach_text
                            else:
                                run.text = f"See {client_name} proposal details."

                        # Apply text replacements
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for old, new in replacements.items():
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)

    # --- Update Slide 3 Key Executive Summary body (shape S84) ---
    if len(prs.slides) >= 3:
        slide3 = prs.slides[2]
        for shape in slide3.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text
            if "Context & Objective" in full_text and "Scope of Solution" in full_text:
                # This is the main exec summary body — update BNI Life references
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)

    safe_name = re.sub(r'[^\w\s-]', '', client_name).strip().replace(' ', '_')
    out_path = os.path.join(OUTPUT_DIR, safe_name, f"Sembuh AI - {client_name} - Proposal v1.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"  Proposal saved: {out_path}")
    return out_path


# ── Google Drive Upload ───────────────────────────────────────

def upload_to_drive(files, client_name):
    """Upload generated files to Google Drive under client folder."""
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), "gmcp"))
    from google_tools import drive_create_file, drive_upload_file

    # Create client folder under main folder
    result = drive_create_file(
        title=f"{client_name}",
        mime_type="folder",
        folder_id=GOOGLE_FOLDER_ID,
    )
    print(f"  Drive folder: {result}")

    # Extract folder ID from result
    folder_id_match = re.search(r"id: ([a-zA-Z0-9_-]+)", result)
    client_folder_id = folder_id_match.group(1) if folder_id_match else GOOGLE_FOLDER_ID

    # Create subfolders
    for subfolder in ["Proposals", "Quotations", "Pricing"]:
        sub_result = drive_create_file(
            title=subfolder,
            mime_type="folder",
            folder_id=client_folder_id,
        )
        sub_id_match = re.search(r"id: ([a-zA-Z0-9_-]+)", sub_result)
        sub_id = sub_id_match.group(1) if sub_id_match else client_folder_id

        # Upload matching file
        for f in files:
            basename = os.path.basename(f).lower()
            if subfolder == "Proposals" and "proposal" in basename:
                drive_upload_file(f, folder_id=sub_id)
                print(f"  Uploaded {os.path.basename(f)} -> {subfolder}/")
            elif subfolder == "Quotations" and "quotation" in basename:
                drive_upload_file(f, folder_id=sub_id)
                print(f"  Uploaded {os.path.basename(f)} -> {subfolder}/")
            elif subfolder == "Pricing" and "pricing" in basename:
                drive_upload_file(f, folder_id=sub_id)
                print(f"  Uploaded {os.path.basename(f)} -> {subfolder}/")

    return client_folder_id


# ── Gmail Draft ───────────────────────────────────────────────

def create_email_draft(deal, files):
    """Create Gmail draft with generated files attached."""
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), "gmcp"))
    from google_tools import gmail_create_draft

    client_name = deal["name"]
    products_str = ", ".join(deal["products"])
    stage = deal["stage"]

    subject = f"Sembuh AI — Proposal for {client_name} ({products_str})"
    body = f"""Dear {client_name} Team,

Thank you for your interest in Sembuh AI's solutions. Please find attached our proposal and quotation for your review.

Proposal Summary:
- Client: {client_name}
- Solutions: {products_str}
- Stage: {stage}
- Monthly Claim Volume: {int(deal['monthly_claim_vol']):,}
- Members Covered: {int(deal['members_covered']):,}

Attached Documents:
1. Proposal — Overview of our solutions and implementation plan
2. Quotation — Detailed pricing breakdown

We look forward to discussing next steps. Please don't hesitate to reach out with any questions.

Best regards,
Sembuh AI Sales Team
www.sembuh.ai

---
This is an auto-generated draft. Please review before sending.
Generated on {TODAY_STR}
"""

    result = gmail_create_draft(
        to="sales@sembuh.ai",
        subject=subject,
        body=body,
        attachments=[f for f in files if "pricing" not in os.path.basename(f).lower()],
    )
    print(f"  Gmail draft: {result}")
    return result


# ── Slack Notification ────────────────────────────────────────

def notify_slack(deals, generated_files):
    """Send summary notification to Slack channel."""
    blocks = [
        {
            "type": "header",
            "text": {"type": "plain_text", "text": "Sales Intelligence — Documents Generated"}
        },
        {
            "type": "section",
            "text": {
                "type": "mrkdwn",
                "text": f"*{len(deals)} deals processed* from Monday.com on {TODAY_STR}"
            }
        },
        {"type": "divider"},
    ]

    for deal, files in zip(deals, generated_files):
        status_emoji = ":large_green_circle:" if deal["stage"] not in ["Lost"] else ":red_circle:"
        products_str = ", ".join(deal["products"])
        file_names = "\n".join([f"  • `{os.path.basename(f)}`" for f in files])

        blocks.append({
            "type": "section",
            "text": {
                "type": "mrkdwn",
                "text": (
                    f"{status_emoji} *{deal['name']}*\n"
                    f"Stage: {deal['stage']} | Value: ${deal['deal_value']:,.0f}\n"
                    f"Products: {products_str}\n"
                    f"Files generated:\n{file_names}"
                )
            }
        })

    blocks.append({"type": "divider"})
    blocks.append({
        "type": "section",
        "text": {
            "type": "mrkdwn",
            "text": (
                ":email: *Gmail drafts created* for all active deals\n"
                ":file_folder: *Files uploaded* to Google Drive\n"
                ":memo: Review drafts before sending to clients"
            )
        }
    })

    # Send via Slack API
    headers = {
        "Authorization": f"Bearer {SLACK_BOT_TOKEN}",
        "Content-Type": "application/json",
    }
    payload = {
        "channel": SLACK_CHANNEL_ID,
        "text": f"Sales Intelligence: {len(deals)} deal documents generated",
        "blocks": blocks,
    }
    resp = requests.post(f"{SLACK_API_URL}/chat.postMessage", json=payload, headers=headers)
    data = resp.json()
    if data.get("ok"):
        print(f"  Slack notification sent to {SLACK_CHANNEL_ID}")
    else:
        print(f"  Slack error: {data.get('error', 'unknown')}")
    return data


# ── Main ──────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("Sales Intelligence — Document Generator")
    print("=" * 60)

    # 1. Fetch deals from Monday.com
    print("\n[1/5] Fetching deals from Monday.com...")
    boards = fetch_boards()
    if not boards:
        print("No boards found in workspace. Exiting.")
        sys.exit(1)

    # Find Deals board (exact match preferred, then substring)
    deals_board = None
    for b in boards:
        if b["name"].strip().lower() == "deals":
            deals_board = b
            break
    if not deals_board:
        for b in boards:
            if "deal" in b["name"].lower() and "subitem" not in b["name"].lower():
                deals_board = b
                break
    if not deals_board:
        deals_board = boards[0]

    print(f"  Using board: {deals_board['name']} (id: {deals_board['id']})")

    raw_deals = fetch_deals(deals_board["id"])
    print(f"  Found {len(raw_deals)} deals")

    deals = [parse_deal(d) for d in raw_deals]

    # 2. Generate documents for each deal
    print("\n[2/5] Generating documents...")
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    all_generated_files = []
    for i, deal in enumerate(deals):
        print(f"\n  --- {deal['name']} ({deal['stage']}) ---")
        files = []

        quotation_path = generate_quotation(deal, seq_num=i + 1)
        files.append(quotation_path)

        pricing_path = generate_pricing_internal(deal, deals)
        files.append(pricing_path)

        proposal_path = generate_proposal(deal)
        files.append(proposal_path)

        all_generated_files.append(files)

    # 3. Upload to Google Drive
    print("\n[3/5] Uploading to Google Drive...")
    for deal, files in zip(deals, all_generated_files):
        print(f"\n  --- {deal['name']} ---")
        try:
            upload_to_drive(files, deal["name"])
        except Exception as e:
            print(f"  Drive upload error: {e}")

    # 4. Create Gmail drafts
    print("\n[4/5] Creating Gmail drafts...")
    for deal, files in zip(deals, all_generated_files):
        if deal["stage"] == "Lost":
            print(f"  Skipping {deal['name']} (Lost)")
            continue
        print(f"\n  --- {deal['name']} ---")
        try:
            create_email_draft(deal, files)
        except Exception as e:
            print(f"  Gmail draft error: {e}")

    # 5. Notify Slack
    print("\n[5/5] Sending Slack notification...")
    try:
        notify_slack(deals, all_generated_files)
    except Exception as e:
        print(f"  Slack notification error: {e}")

    print("\n" + "=" * 60)
    print("Done! All documents generated.")
    print(f"Output directory: {OUTPUT_DIR}")
    print("=" * 60)


if __name__ == "__main__":
    main()
